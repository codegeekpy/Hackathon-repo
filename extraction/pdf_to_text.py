import os
import logging
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def extract_text_from_docx(file_path):
    try:
        document = Document(file_path)
        full_text = []
        for para in document.paragraphs:
            full_text.append(para.text)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)
        logger.info(f"Successfully extracted text from DOCX: {file_path}")
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting from DOCX '{file_path}': {e}", exc_info=True)
        return f"Error extracting from DOCX: {e}"

def extract_text_from_pdf(file_path):
    full_text = []
    try:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                full_text.append(page_text)
            else:
                logger.warning(f"No text extracted from page {i+1} of PDF: {file_path}. It might be an image-based page or have no text content.")
        logger.info(f"Successfully extracted text from PDF: {file_path}")
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting from PDF '{file_path}': {e}", exc_info=True)
        return f"Error extracting from PDF: {e}"

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide_num, slide in enumerate(prs.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    if shape.text:
                        full_text.append(shape.text)
                    else:
                        logger.debug(f"Empty text content in shape {shape_num+1} on slide {slide_num+1} of PPTX: {file_path}")
        logger.info(f"Successfully extracted text from PPTX: {file_path}")
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting from PPTX '{file_path}': {e}", exc_info=True)
        return f"Error extracting from PPTX: {e}"

def extract_text_from_document(file_path):
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        return f"File not found: {file_path}"

    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".docx":
        return extract_text_from_docx(file_path)
    elif file_extension == ".pdf":
        return extract_text_from_pdf(file_path)
    elif file_extension == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        logger.warning(f"Unsupported file type encountered: {file_extension} for file {file_path}")
        return f"Unsupported file type: {file_extension}"

if __name__ == "__main__":
    # Remove or comment out the dummy file creation if you only want to process user-provided files
    # doc = Document()
    # ... (dummy file creation code) ...

    while True:
        file_path = input("Enter the full path to the document,if u selected copy as path please remove the double codes (or 'q' to quit): ")
        if file_path.lower() == 'q':
            break

        if not file_path:
            logger.warning("No file path entered. Please try again.")
            continue

        print(f"\n--- Processing {file_path} ---")
        extracted_text = extract_text_from_document(file_path)
        print(f"\nExtracted Text from {file_path}:\n{extracted_text}")

    logger.info("Document extraction process finished.")